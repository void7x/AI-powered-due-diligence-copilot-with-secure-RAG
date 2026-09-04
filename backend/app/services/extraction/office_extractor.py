"""DOCX and PPTX extraction with heading/slide-aware structure."""
from __future__ import annotations

import math

from app.core.errors import MalformedDocumentError
from app.services.extraction.base import ExtractedPage, ExtractionResult

_APPROX_PAGE_CHARS = 2400


def _approx_pages(blocks: list[tuple[str | None, str]]) -> list[ExtractedPage]:
    """DOCX has no real pages: pack blocks into ~page-sized pseudo pages."""
    pages: list[ExtractedPage] = []
    buf: list[str] = []
    page_no = 1
    for style, text in blocks:
        buf.append(text)
        if sum(len(b) for b in buf) >= _APPROX_PAGE_CHARS:
            pages.append(ExtractedPage(page_number=page_no, text="\n\n".join(buf),
                                       meta={"char_count": sum(len(b) for b in buf)},
                                       approximate_page=True))
            page_no += 1
            buf = []
    if buf:
        pages.append(ExtractedPage(page_number=page_no, text="\n\n".join(buf),
                                   meta={"char_count": sum(len(b) for b in buf)},
                                   approximate_page=True))
    return pages or [ExtractedPage(page_number=1, text="")]


def extract_docx(path: str) -> ExtractionResult:
    try:
        import docx  # python-docx
        document = docx.Document(path)
    except Exception as exc:  # noqa: BLE001
        raise MalformedDocumentError(f"Could not open DOCX: {exc}") from exc
    blocks: list[tuple[str | None, str]] = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style is not None else ""
        prefix = f"## {text}" if style.startswith("Heading") else text
        blocks.append((style, prefix))
    for table in document.tables:  # type: ignore[attr-defined]
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if rows:
            blocks.append((None, "[TABLE]\n" + "\n".join(" | ".join(r) for r in rows)))
    if not blocks:
        raise MalformedDocumentError("DOCX contains no extractable text.")
    return ExtractionResult(pages=_approx_pages(blocks), parser="python-docx", page_basis="approximate")


def extract_pptx(path: str) -> ExtractionResult:
    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception as exc:  # noqa: BLE001
        raise MalformedDocumentError(f"Could not open PPTX: {exc}") from exc
    pages: list[ExtractedPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(("# " if para.level == 0 and shape == slide.shapes.title else "") + text)
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]  # type: ignore[attr-defined]
                lines.append("[TABLE]\n" + "\n".join(" | ".join(r) for r in rows))
        if lines:
            pages.append(ExtractedPage(page_number=i, text="\n".join(lines),
                                       meta={"char_count": sum(len(l) for l in lines)}))
    if not pages:
        raise MalformedDocumentError("PPTX contains no extractable text.")
    return ExtractionResult(pages=pages, parser="python-pptx", page_basis="slide")
